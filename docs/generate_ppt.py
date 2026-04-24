# -*- coding: utf-8 -*-
"""
毕业答辩PPT生成脚本
基于项目Wiki和论文内容生成答辩演示文稿
要求：字号24以上，演讲时长5-7分钟（约15-20页）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 颜色配置
# ============================================================
COLOR_PRIMARY = RGBColor(0x1A, 0x5F, 0x9E)      # 深蓝 - 主色调
COLOR_SECONDARY = RGBColor(0x2E, 0x8B, 0x57)    # 深绿 - 辅助色
COLOR_ACCENT = RGBColor(0xE6, 0x6C, 0x00)       # 橙色 - 强调色
COLOR_DARK = RGBColor(0x2C, 0x3E, 0x50)         # 深色文字
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_SUBTITLE = RGBColor(0x5D, 0x6D, 0x7E)

# ============================================================
# 字号配置（全部>=24pt）
# ============================================================
FONT_SIZE_TITLE = Pt(44)        # 封面标题
FONT_SIZE_SUBTITLE = Pt(28)     # 副标题
FONT_SIZE_SECTION = Pt(36)      # 章节标题
FONT_SIZE_HEADING = Pt(32)      # 页面标题
FONT_SIZE_BODY = Pt(24)         # 正文（最小24pt）
FONT_SIZE_SMALL = Pt(24)        # 小字（最小24pt）

# ============================================================
# 幻灯片尺寸（16:9）
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size, bold=True, color=COLOR_DARK, align=PP_ALIGN.LEFT):
    """添加标题文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "微软雅黑"
    p.alignment = align
    return shape


def add_body_text(slide, lines, left, top, width, height, font_size=FONT_SIZE_BODY, color=COLOR_DARK, bullet=True, line_spacing=1.5):
    """添加正文文本框，支持多行"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(font_size.pt * (line_spacing - 1))
        if bullet and line.strip():
            p.level = 0
    return shape


def add_bullet_text(slide, items, left, top, width, height, font_size=FONT_SIZE_BODY, color=COLOR_DARK, line_spacing=1.5):
    """添加带项目符号的文本"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {item}"
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(font_size.pt * (line_spacing - 1))
    return shape


def add_decorated_title(slide, title, subtitle=None):
    """添加带装饰的页面标题栏"""
    # 顶部装饰条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # 标题
    add_title_shape(slide, title, Inches(0.6), Inches(0.5), Inches(12), Inches(0.9),
                    FONT_SIZE_HEADING, bold=True, color=COLOR_PRIMARY)

    # 副标题
    if subtitle:
        add_title_shape(slide, subtitle, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
                        FONT_SIZE_SMALL, bold=False, color=COLOR_SUBTITLE)
    return Inches(1.7)


def create_presentation():
    """创建答辩PPT"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ============================================================
    # 第1页：封面
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, COLOR_WHITE)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = COLOR_PRIMARY; top_bar.line.fill.background()

    # 学校/学院信息
    add_title_shape(slide, "本科毕业论文（设计）答辩", Inches(0), Inches(1.5), SLIDE_WIDTH, Inches(0.8),
                    Pt(32), bold=True, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)

    # 主标题
    add_title_shape(slide, "多模态智能校园助手系统的设计与实现", Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(1.0),
                    FONT_SIZE_TITLE, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    # 英文副标题
    add_title_shape(slide, "Design and Implementation of Multimodal Intelligent Campus Assistant System",
                    Inches(0), Inches(3.9), SLIDE_WIDTH, Inches(0.6),
                    Pt(20), bold=False, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)

    # 底部信息
    info_lines = [
        "答辩人：XXX",
        "学  号：2227010250",
        "专  业：计算机科学与技术",
        "指导教师：XXX 教授"
    ]
    info_text = "\n".join(info_lines)
    add_title_shape(slide, info_text, Inches(0), Inches(5.3), SLIDE_WIDTH, Inches(1.8),
                    FONT_SIZE_BODY, bold=False, color=COLOR_DARK, align=PP_ALIGN.CENTER)

    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15))
    bottom_bar.fill.solid(); bottom_bar.fill.fore_color.rgb = COLOR_PRIMARY; bottom_bar.line.fill.background()

    # ============================================================
    # 第2页：目录
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "目  录")

    contents = [
        "01  研究背景与意义",
        "02  系统总体架构",
        "03  核心功能模块",
        "04  技术创新与亮点",
        "05  总结与展望"
    ]
    for i, item in enumerate(contents):
        y = top_y + Inches(0.3 + i * 0.85)
        # 序号圆形背景
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.6), Inches(0.6))
        circle.fill.solid(); circle.fill.fore_color.rgb = COLOR_PRIMARY; circle.line.fill.background()
        num_tf = circle.text_frame; num_tf.paragraphs[0].text = item[:2]
        num_tf.paragraphs[0].font.size = Pt(22); num_tf.paragraphs[0].font.bold = True
        num_tf.paragraphs[0].font.color.rgb = COLOR_WHITE; num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_title_shape(slide, item[4:], Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.6),
                        FONT_SIZE_BODY, bold=True, color=COLOR_DARK)

    # ============================================================
    # 第3页：研究背景
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "01  研究背景与意义", "Research Background and Significance")

    items = [
        "高校信息化正从\"数字校园\"向\"智慧校园\"跃迁，大语言模型与多模态技术日趋成熟",
        "现有校园系统功能分散、入口独立，师生查询课表、成绩、办事指南时效率低下",
        "传统系统交互方式单一，不支持语音、图像等多模态输入，用户体验不佳",
        "教务数据涉及隐私，缺少严格的身份校验与越权约束，存在数据泄露风险"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.2), Inches(11.5), Inches(3.5))

    # 痛点总结框
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.0), Inches(11.5), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "核心痛点：入口分散  ·  交互单一  ·  隐私薄弱  ·  响应滞后"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第4页：研究意义
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "01  研究意义", "Research Significance")

    left_items = [
        "统一对话入口：整合教务查询、办事指南、规章制度等分散信息源",
        "多模态交互：支持文本、语音、图像输入，实现\"所见即所问\"",
        "智能意图路由：自动识别查询类型，精准匹配结构化查询或知识检索",
        "隐私安全保障：JWT认证 + 行级数据隔离 + DID脱敏，确保数据安全"
    ]
    add_bullet_text(slide, left_items, Inches(0.8), top_y + Inches(0.2), Inches(11.5), Inches(3.5))

    # 价值框
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9); box.line.color.rgb = COLOR_SECONDARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "应用价值：为学生提供24小时在线的\"私人助理\"服务，显著降低信息获取门槛"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_SECONDARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第5页：系统总体架构
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "02  系统总体架构", "System Architecture")

    # 架构分层展示
    layers = [
        ("表现层", "Vue 3 + Vite + Pinia + Tailwind CSS", COLOR_PRIMARY),
        ("业务层", "FastAPI + 认证/查询/管理路由 + 服务编排", RGBColor(0x27, 0xAE, 0x60)),
        ("数据层", "MySQL 8.0 + SQLAlchemy AsyncIO ORM", RGBColor(0x8E, 0x44, 0xAD)),
        ("外部服务", "DashScope(Qwen) + 百炼检索API + Redis缓存", RGBColor(0xD3, 0x54, 0x00)),
    ]
    for i, (name, desc, color) in enumerate(layers):
        y = top_y + Inches(0.3 + i * 1.0)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.75))
        rect.fill.solid(); rect.fill.fore_color.rgb = color; rect.line.fill.background()
        tf = rect.text_frame
        tf.paragraphs[0].text = f"{name}：{desc}"
        tf.paragraphs[0].font.size = FONT_SIZE_BODY; tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第6页：核心业务流程
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "02  核心业务流程", "Core Business Process")

    flow_steps = [
        "① 多模态输入：文本 / 图像(Qwen-VL) / 语音(ASR) → 统一查询文本",
        "② 缓存优先：Redis 缓存命中即快速返回，降低延迟与成本",
        "③ 安全检查：并发执行危险内容检测 + 隐私越权检测",
        "④ 意图分类：LLM 识别 structured / vector / hybrid / smalltalk",
        "⑤ 查询执行：结构化SQL查询 或 向量检索(RAG) 或 混合执行",
        "⑥ LLM总结：生成自然语言回答，SSE流式推送至前端"
    ]
    add_bullet_text(slide, flow_steps, Inches(0.8), top_y + Inches(0.2), Inches(11.5), Inches(4.5),
                    line_spacing=1.3)

    # ============================================================
    # 第7页：核心功能 - 多模态输入处理
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 多模态输入处理", "Multimodal Input Processing")

    items = [
        "图像理解：支持 JPEG/PNG，自动缩放压缩，调用 Qwen-VL 生成图片描述",
        "语音识别：支持 WAV/MP3，ffmpeg 转 16kHz 单声道后调用 Paraformer ASR",
        "文本拼接：将图像/语音转文本与原始文本合并，形成统一查询语句",
        "图片纯问答：针对\"解释/分析\"类问题直接基于图片描述回答，避免无关检索"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # 示意图
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "输入融合流程：文本 + [图片内容] + [语音转文字] → 统一查询文本 → 后续处理"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第8页：核心功能 - 意图识别与查询重写
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 意图识别与查询重写", "Intent Classification & Query Rewrite")

    items = [
        "意图分类：基于 Qwen 模型将查询归类为 structured / vector / hybrid / smalltalk",
        "查询重写：结合最近 3~5 轮会话历史，补齐代词与缩写的缺失信息",
        "上下文截断：针对 LLM 输入长度限制进行消息裁剪，避免超限失败",
        "动态修正：执行后根据实际返回内容修正意图类型，提升准确率"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # 示例框
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0); box.line.color.rgb = COLOR_ACCENT
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "示例：\"那这门课呢？\" → 重写为 \"张三的软件工程课什么时候上？\""
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_ACCENT
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第9页：核心功能 - 结构化数据查询
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 结构化数据查询", "Structured Data Query (Text2SQL)")

    items = [
        "隐私约束：所有查询严格限定当前登录学生身份，防止越权访问他人数据",
        "工具规划：根据意图选择成绩查询、课表查询、选课查询、教师通讯录等工具",
        "自动学期解析：根据系统日期智能推断当前/上/下学期，支持相对时间表达",
        "结果美化：字段名翻译、学期ID格式化、布尔值人性化，提升可读性"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # 支持查询类型
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "支持查询：成绩 · 课表 · 选课 · 个人信息 · 教师通讯录 · 院系专业目录"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第10页：核心功能 - 知识库检索(RAG)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 知识库检索 (RAG)", "Retrieval-Augmented Generation")

    items = [
        "百炼检索：封装阿里云百炼检索 API 为 LangChain Retriever，支持异步召回",
        "查询拆解：将复杂问题拆分为关键词短语，提升知识库召回率与相关性",
        "重排融合：检索结果与结构化数据去重、筛选、重排，构建高密度上下文",
        "上下文截断：对大体量检索结果进行头尾截断，保留关键信息片段"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9); box.line.color.rgb = COLOR_SECONDARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "应用场景：办事流程 · 规章制度 · 奖学金申请 · 请假/证明办理指南"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_SECONDARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第11页：核心功能 - SSE流式响应与缓存
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 流式响应与缓存优化", "Streaming Response & Cache Optimization")

    left_items = [
        "SSE 流式输出：使用 StreamingResponse 逐块推送 Token，前端逐字渲染",
        "结束元数据：最后一条数据携带意图类型、响应耗时、缓存状态等信息",
        "敏感度分级缓存：敏感查询短 TTL，普通查询长 TTL，平衡效率与准确性",
        "版本控制：日期敏感按日桶失效，课表敏感按版本号失效，防止陈旧结果"
    ]
    add_bullet_text(slide, left_items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.5), Inches(11.5), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "缓存键设计：chat_cache:{version}:{did}:{md5(query)}"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第12页：核心功能 - 安全与隐私保护
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "03  核心功能 — 安全与隐私保护", "Security & Privacy Protection")

    items = [
        "JWT 认证：学生与管理员双通道鉴权，AES-CBC 加密密码传输",
        "危险内容检测：LLM 识别自杀/暴力倾向，触发危机干预提示",
        "隐私越权拦截：禁止查询他人学号，违者返回隐私保护提示",
        "DID 脱敏：基于 student_id 生成稳定单向 DID，会话历史隔离存储",
        "降级策略：Redis/LLM/数据库异常时自动降级，确保核心链路可用"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(4.0))

    # ============================================================
    # 第13页：技术创新与亮点（一）
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "04  技术创新与亮点（一）", "Technical Innovation & Highlights")

    items = [
        "LangChain LCEL 编排：构建高扩展性的并发调用链路，支持同步/流式双模输出",
        "混合查询执行：结构化 SQL + 向量检索 RAG 动态融合，覆盖教务数据与校园知识",
        "多轮对话上下文重写：结合会话历史补齐信息，解决代词与缩写的理解难题",
        "受控工具调用：模型仅生成调用计划，后端严格绑定当前身份执行查询，保障安全"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # 亮点标签
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.8), Inches(11.5), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9); box.line.color.rgb = COLOR_SECONDARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "核心优势：多模态统一  ·  智能路由  ·  安全可控  ·  流式体验"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_SECONDARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第14页：技术创新与亮点（二）
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "04  技术创新与亮点（二）", "Technical Innovation & Highlights")

    items = [
        "并发优化：安全检查与查询重写并行执行，缩短端到端响应延迟",
        "连接池管理：流式阶段提前回滚数据库会话，避免长连接占用",
        "敏感缓存分级：按关键词判断敏感度，不同 TTL 策略精细化控制",
        "Docker 容器化：全套环境通过 Docker Compose 统一部署，轻量且环境一致"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # 技术栈标签
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.8), Inches(11.5), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "技术栈：Vue 3 + FastAPI + SQLAlchemy + Redis + DashScope + LangChain + Docker"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第15页：数据库设计
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "04  数据库设计", "Database Design")

    items = [
        "核心实体：学生、课程、成绩、课表、教师、教室、院系、专业、班级等",
        "关系约束：一对多/多对多外键约束，索引优化查询性能",
        "聊天日志：DID 脱敏关联、sender 类型、response_time_ms 等审计字段",
        "调课管理：Schedule_Adjustment 支持调课审批、版本控制与冲突快照"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.8), Inches(11.5), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0); box.line.color.rgb = COLOR_ACCENT
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "设计原则：规范化建模  ·  隐私脱敏  ·  审计可追溯  ·  版本化管理"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_ACCENT
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第16页：系统展示 - 前端界面
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "04  系统展示 — 前端界面", "System Frontend")

    items = [
        "左右分栏布局：左侧导航功能区，右侧对话主交互区",
        "多模态输入：支持文本输入、图片上传、语音录制与播放",
        "流式渲染：SSE 逐字打印效果，实时显示意图、耗时、缓存状态",
        "会话管理：新建对话、搜索历史、会话切换与本地持久化"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(3.8), Inches(11.5), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = COLOR_LIGHT_BG; box.line.color.rgb = COLOR_PRIMARY
    box_tf = box.text_frame; box_tf.word_wrap = True
    box_p = box_tf.paragraphs[0]
    box_p.text = "交互亮点：低门槛上手  ·  实时反馈  ·  多模态融合  ·  上下文连贯"
    box_p.font.size = FONT_SIZE_BODY; box_p.font.bold = True; box_p.font.color.rgb = COLOR_PRIMARY
    box_p.alignment = PP_ALIGN.CENTER

    # ============================================================
    # 第17页：总结
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "05  总结与展望", "Summary & Future Work")

    items = [
        "系统实现了多模态输入、智能意图路由、结构化查询与 RAG 融合的校园问答服务",
        "通过 JWT 认证、隐私脱敏、越权检测等机制构建了完整的安全防护链路",
        "SSE 流式输出与 Redis 缓存优化显著提升了用户体验与系统性能",
        "容器化部署方案使系统具备良好的可移植性与扩展性"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(3.5))

    # ============================================================
    # 第18页：未来展望
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)
    top_y = add_decorated_title(slide, "05  未来展望", "Future Work")

    items = [
        "引入对话摘要与长期记忆机制，支持跨会话的上下文保持",
        "扩展更多检索器与重排策略，提升知识库召回准确率",
        "支持更多媒体格式（视频、文档），丰富多模态交互能力",
        "完善监控与可观测性指标，构建运维 Dashboard",
        "推广至更多高校场景，形成可复制的智慧校园服务范式"
    ]
    add_bullet_text(slide, items, Inches(0.8), top_y + Inches(0.3), Inches(11.5), Inches(4.0))

    # ============================================================
    # 第19页：致谢
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)

    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = COLOR_PRIMARY; top_bar.line.fill.background()

    add_title_shape(slide, "致  谢", Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.0),
                    Pt(48), bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    add_title_shape(slide, "感谢各位老师的指导与评审！",
                    Inches(0), Inches(3.8), SLIDE_WIDTH, Inches(0.8),
                    FONT_SIZE_HEADING, bold=False, color=COLOR_SUBTITLE, align=PP_ALIGN.CENTER)

    add_title_shape(slide, "恳请各位老师批评指正",
                    Inches(0), Inches(4.6), SLIDE_WIDTH, Inches(0.8),
                    FONT_SIZE_BODY, bold=False, color=COLOR_DARK, align=PP_ALIGN.CENTER)

    # 底部装饰条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15))
    bottom_bar.fill.solid(); bottom_bar.fill.fore_color.rgb = COLOR_PRIMARY; bottom_bar.line.fill.background()

    # ============================================================
    # 保存PPT
    # ============================================================
    output_path = r"D:\Win\毕设\pro\ai-aiss-main\docs\毕业答辩PPT.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    create_presentation()
